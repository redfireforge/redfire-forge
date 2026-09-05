#!/usr/bin/env python3
"""Generate a RedfireForge training PowerPoint deck."""

import json
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BRAND_DARK = RGBColor(0x0F, 0x17, 0x2A)
BRAND_BLUE = RGBColor(0x3B, 0x5B, 0xDB)
BRAND_LIGHT = RGBColor(0xE8, 0xEB, 0xF5)
BRAND_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BRAND_ORANGE = RGBColor(0xFF, 0x6B, 0x35)
BRAND_GRAY = RGBColor(0x8B, 0x8D, 0x97)
BRAND_GREEN = RGBColor(0x40, 0xC0, 0x57)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=BRAND_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide(slide, items, start_top=Inches(2), left=Inches(1),
                     font_size=20, color=BRAND_WHITE, spacing=Inches(0.55)):
    for i, item in enumerate(items):
        add_text_box(slide, left, start_top + spacing * i,
                     Inches(11), Inches(0.5), f"•  {item}",
                     font_size=font_size, color=color)


def add_two_column(slide, left_items, right_items, top=Inches(2.2)):
    for i, item in enumerate(left_items):
        add_text_box(slide, Inches(0.8), top + Inches(0.5) * i,
                     Inches(5.5), Inches(0.45), f"•  {item}",
                     font_size=18, color=BRAND_WHITE)
    for i, item in enumerate(right_items):
        add_text_box(slide, Inches(7), top + Inches(0.5) * i,
                     Inches(5.5), Inches(0.45), f"•  {item}",
                     font_size=18, color=BRAND_WHITE)


def add_section_title(slide, title, subtitle=""):
    set_slide_bg(slide, BRAND_BLUE)
    add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.2),
                 title, font_size=44, bold=True, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
                     subtitle, font_size=22, color=BRAND_LIGHT,
                     alignment=PP_ALIGN.CENTER)


def add_slide_title(slide, title, subtitle=""):
    set_slide_bg(slide, BRAND_DARK)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x15, 0x1F, 0x38)
    shape.line.fill.background()
    add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
                 title, font_size=32, bold=True, color=BRAND_WHITE)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.4),
                     subtitle, font_size=16, color=BRAND_GRAY)


def add_table_slide(slide, title, headers, rows, subtitle=""):
    add_slide_title(slide, title, subtitle)
    cols = len(headers)
    table_rows = len(rows) + 1
    left = Inches(0.8)
    top = Inches(1.8)
    width = Inches(11.5)
    row_height = Inches(0.5)
    tbl = slide.shapes.add_table(
        table_rows, cols, left, top, width, row_height * table_rows).table

    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = BRAND_WHITE
            p.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_BLUE

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r + 1, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.color.rgb = BRAND_WHITE
                p.font.name = "Calibri"
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                RGBColor(0x1A, 0x25, 0x40) if r % 2 == 0
                else RGBColor(0x15, 0x1F, 0x38))


def generate(version: str, output_dir: str):
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank = prs.slide_layouts[6]

    # ── Slide 1: Title ──
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, BRAND_DARK)
    add_text_box(s, Inches(1), Inches(1.8), Inches(11), Inches(1.5),
                 "🔥 RedfireForge", font_size=60, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
                 "API Performance Studio", font_size=28,
                 color=BRAND_LIGHT, alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
                 "Fire. Measure. Validate.", font_size=22,
                 color=BRAND_ORANGE, alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                 f"Training Guide — v{version}", font_size=18,
                 color=BRAND_GRAY, alignment=PP_ALIGN.CENTER)

    # ── Slide 2: Agenda ──
    s = prs.slides.add_slide(blank)
    add_slide_title(s, "Agenda")
    items = [
        "1.  What is RedfireForge?",
        "2.  Installation & Setup",
        "3.  Core Concepts: Environments, Microservices, Feature Groups",
        "4.  Creating & Organizing Tests",
        "5.  Authentication Configuration",
        "6.  Running Performance Tests",
        "7.  Analyzing Results",
        "8.  Export & Import",
        "9.  Desktop vs Web Mode",
        "10. Tips & Best Practices",
    ]
    add_bullet_slide(s, items, start_top=Inches(1.8), font_size=22,
                     spacing=Inches(0.52))

    # ── Slide 3: What is RedfireForge ──
    s = prs.slides.add_slide(blank)
    add_section_title(s, "What is RedfireForge?",
                      "A visual API performance testing tool")

    # ── Slide 4: Overview ──
    s = prs.slides.add_slide(blank)
    add_slide_title(s, "What is RedfireForge?", "Overview")
    items = [
        "Desktop & web application for API performance testing",
        "Define HTTP tests visually — no scripting required",
        "Execute with configurable concurrency (sequential, parallel, ramp-up)",
        "Validate responses with a visual JSON builder",
        "Analyze results with built-in dashboards and charts",
        "Built with React + TypeScript + Vite + Tauri",
    ]
    add_bullet_slide(s, items, start_top=Inches(1.8))

    # ── Slide 5: Key Features ──
    s = prs.slides.add_slide(blank)
    add_slide_title(s, "Key Features")
    add_two_column(s,
        ["Hierarchical test organization",
         "Drag-and-drop reordering",
         "OAuth2 authentication",
         "Global auth profiles",
         "Environment management",
         "Microservice filtering"],
        ["Configurable concurrency",
         "Response validation (JSON)",
         "Results dashboard & charts",
         "Export / Import with conflict resolution",
         "Desktop app (macOS, Windows, Linux)",
         "Cross-platform data portability"])

    # ── Slide 6: Installation ──
    s = prs.slides.add_slide(blank)
    add_section_title(s, "Installation & Setup")

    # ── Slide 7: Installation Details ──
    s = prs.slides.add_slide(blank)
    add_slide_title(s, "Installation", "Desktop Application")
    items = [
        "macOS: Open the .dmg file → drag RedfireForge to Applications",
        "Windows: Run the .msi installer → follow the wizard",
        "Linux: Use the .AppImage (make executable) or install the .deb",
        "",
        "First-time macOS: Right-click → Open → click \"Open\" to bypass Gatekeeper",
        "",
        "Web mode: Open http://localhost:5173 in browser (dev server required)",
    ]
    add_bullet_slide(s, items, start_top=Inches(1.8), spacing=Inches(0.5))

    # ── Slide 8: Core Concepts Section ──
    s = prs.slides.add_slide(blank)
    add_section_title(s, "Core Concepts",
                      "Environments, Microservices & Feature Groups")

    # ── Slide 9: Hierarchy ──
    s = prs.slides.add_slide(blank)
    add_slide_title(s, "Test Hierarchy", "How tests are organized")
    items = [
        "Environment  →  dev, test, staging, prod, etc.",
        "Microservice  →  order-api, user-service, etc.",
        "Feature Group  →  A collection of related test scenarios",
        "    └─ Scenario  →  A specific test flow (e.g., \"Onboarding Flow\")",
        "        └─ Test  →  An individual HTTP request with assertions",
    ]
    add_bullet_slide(s, items, start_top=Inches(1.8), font_size=22,
                     spacing=Inches(0.65))

    # ── Slide 10: Environments & Microservices ──
    s = prs.slides.add_slide(blank)
    add_slide_title(s, "Environments & Microservices",
                    "Configure in Settings → Sidebar")
    items = [
        "Settings → manage Environments (dev, test, staging, prod, ...)",
        "Settings → manage Microservices",
        "Sidebar filters Feature Groups by selected Environment + Microservice",
        "Each Feature Group is tagged with an Environment and Microservice",
    ]
    add_bullet_slide(s, items, start_top=Inches(1.8), font_size=20,
                     spacing=Inches(0.65))

    # ── Slide 11: Creating Tests Section ──
    s = prs.slides.add_slide(blank)
    add_section_title(s, "Creating & Organizing Tests")

    # ── Slide 12: Feature Groups ──
    s = prs.slides.add_slide(blank)
    add_slide_title(s, "Feature Groups",
                    "Container for related scenarios")
    items = [
        "Click \"+ Add Feature Group\" in the Feature Groups tab",
        "Assign an Environment and Microservice",
        "Add Scenarios inside each Feature Group",
        "Add Tests inside each Scenario",
        "Use drag-and-drop to reorder or move scenarios between groups",
    ]
    add_bullet_slide(s, items, start_top=Inches(1.8), font_size=20,
                     spacing=Inches(0.6))

    # ── Slide 13: Test Configuration ──
    s = prs.slides.add_slide(blank)
    add_slide_title(s, "Test Configuration", "Setting up an HTTP request")
    add_table_slide.__wrapped__ = True  # skip, use manual
    add_two_column(s,
        ["Name & description",
         "HTTP method (GET, POST, PUT, DELETE, PATCH)",
         "URL (supports environment variables)",
         "Headers (key-value pairs)",
         "Request body (JSON)"],
        ["Expected status code",
         "Response validation rules",
         "Authentication override",
         "Timeout settings",
         "Pre/post request delays"],
        top=Inches(1.8))

    # ── Slide 14: Auth Section ──
    s = prs.slides.add_slide(blank)
    add_section_title(s, "Authentication",
                      "Global Profiles & Inheritance")

    # ── Slide 15: Auth Details ──
    s = prs.slides.add_slide(blank)
    add_slide_title(s, "Authentication Configuration",
                    "4-tier inheritance model")
    items = [
        "Global Auth Profile  →  defined in Settings (e.g., dev-oauth, prod-oauth)",
        "Feature Group Auth  →  inherits from Global, can override",
        "Scenario Auth  →  inherits from Feature Group, can override",
        "Test Auth  →  inherits from Scenario, can override",
        "",
        "Supports: OAuth2 (Client Credentials), Bearer Token, Basic Auth, API Key",
        "Use \"Verify\" button to test credentials before running tests",
    ]
    add_bullet_slide(s, items, start_top=Inches(1.8), font_size=20,
                     spacing=Inches(0.55))

    # ── Slide 16: Test Runner Section ──
    s = prs.slides.add_slide(blank)
    add_section_title(s, "Running Performance Tests",
                      "Test Runner Configuration")

    # ── Slide 17: Test Runner ──
    s = prs.slides.add_slide(blank)
    add_slide_title(s, "Test Runner", "Configure and execute")
    items = [
        "Select a Feature Group, Scenario, or individual Test to run",
        "Configure number of iterations (runs)",
        "Choose execution mode:",
        "    Sequential  →  one request at a time",
        "    Parallel  →  all requests at once",
        "    Ramp-up  →  gradually increase concurrency",
        "View real-time progress during execution",
        "Results auto-save after completion",
    ]
    add_bullet_slide(s, items, start_top=Inches(1.8), font_size=20,
                     spacing=Inches(0.5))

    # ── Slide 18: Results Section ──
    s = prs.slides.add_slide(blank)
    add_section_title(s, "Analyzing Results",
                      "Results Dashboard")

    # ── Slide 19: Results Dashboard ──
    s = prs.slides.add_slide(blank)
    add_slide_title(s, "Results Dashboard", "Review past test runs")
    items = [
        "View all historical test runs with timestamps",
        "Metrics: avg / min / max / p95 / p99 response times",
        "Error rate and status code distribution",
        "Compare multiple runs side by side",
        "Delete old runs to free storage",
        "Export results as JSON or CSV",
    ]
    add_bullet_slide(s, items, start_top=Inches(1.8), font_size=20,
                     spacing=Inches(0.6))

    # ── Slide 20: Export/Import Section ──
    s = prs.slides.add_slide(blank)
    add_section_title(s, "Export & Import",
                      "Share configurations across teams")

    # ── Slide 21: Export ──
    s = prs.slides.add_slide(blank)
    add_slide_title(s, "Export Center", "Settings → Export")
    items = [
        "Export everything or select specific items",
        "Includes: Environments, Microservices, Auth Profiles, Feature Groups, Results",
        "Standardized filename: {env}-{microservice}-{level}-{name}-{timestamp}.json",
        "Desktop default directory: Documents/RedfireForge/",
        "Share .json files with teammates for consistent test setups",
    ]
    add_bullet_slide(s, items, start_top=Inches(1.8), font_size=20,
                     spacing=Inches(0.6))

    # ── Slide 22: Import ──
    s = prs.slides.add_slide(blank)
    add_slide_title(s, "Import Center", "Settings → Import")
    items = [
        "Import previously exported .json files",
        "Automatic conflict detection (ID + name matching)",
        "Per-item resolution: Skip, Overwrite, or Keep Both",
        "Bulk actions: select all / deselect all per section",
        "Preview all items before importing",
        "Desktop opens Documents/RedfireForge/ by default",
    ]
    add_bullet_slide(s, items, start_top=Inches(1.8), font_size=20,
                     spacing=Inches(0.6))

    # ── Slide 23: Desktop vs Web ──
    s = prs.slides.add_slide(blank)
    add_section_title(s, "Desktop vs Web Mode")

    # ── Slide 24: Comparison Table ──
    s = prs.slides.add_slide(blank)
    add_table_slide(s, "Desktop vs Web", 
        ["Feature", "Desktop (Tauri)", "Web (Browser)"],
        [
            ["Storage", "OS AppData (files)", "localStorage"],
            ["HTTP Requests", "Native (no CORS)", "Vite proxy required"],
            ["File Dialogs", "Native OS dialogs", "Browser file picker"],
            ["Default Directory", "Documents/RedfireForge/", "Browser default"],
            ["Installation", ".dmg / .msi / .AppImage", "npm run dev"],
            ["Offline Use", "Yes", "No (needs dev server)"],
            ["Cross-browser Sync", "N/A (single app)", "Not supported"],
        ])

    # ── Slide 25: Tips ──
    s = prs.slides.add_slide(blank)
    add_section_title(s, "Tips & Best Practices")

    # ── Slide 26: Tips Details ──
    s = prs.slides.add_slide(blank)
    add_slide_title(s, "Tips & Best Practices")
    items = [
        "Organize by environment: one Feature Group per env/microservice combo",
        "Use Global Auth Profiles to avoid duplicating credentials",
        "Export your config regularly as backup (Documents/RedfireForge/)",
        "Use drag-and-drop to restructure tests as requirements evolve",
        "Start with Sequential mode, then switch to Parallel for load testing",
        "Check the Results Dashboard after each run — look for p95/p99 spikes",
        "Use the desktop app for best performance and no CORS issues",
    ]
    add_bullet_slide(s, items, start_top=Inches(1.8), font_size=20,
                     spacing=Inches(0.55))

    # ── Slide 27: Quick Reference ──
    s = prs.slides.add_slide(blank)
    add_table_slide(s, "Quick Reference — Keyboard & Navigation",
        ["Action", "Where", "How"],
        [
            ["Switch tabs", "Header bar", "Click Feature Groups / Test Runner / Results"],
            ["Filter tests", "Sidebar", "Select Environment + Microservice"],
            ["Create Feature Group", "Feature Groups tab", "Click \"+ Add Feature Group\""],
            ["Add Scenario", "Inside Feature Group", "Click \"+ Add Scenario\""],
            ["Add Test", "Inside Scenario", "Click \"+ Add Test\""],
            ["Reorder items", "Feature Groups tab", "Drag and drop"],
            ["Open Settings", "Sidebar bottom", "Click ⚙ Settings"],
            ["Export / Import", "Settings modal", "Export Center / Import Center buttons"],
        ])

    # ── Slide 28: End ──
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, BRAND_DARK)
    add_text_box(s, Inches(1), Inches(2), Inches(11), Inches(1.2),
                 "🔥 Ready to Forge!", font_size=52, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
                 "Start testing your APIs with RedfireForge",
                 font_size=24, color=BRAND_LIGHT, alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(1), Inches(5), Inches(11), Inches(0.5),
                 f"v{version} — API Performance Studio",
                 font_size=18, color=BRAND_GRAY, alignment=PP_ALIGN.CENTER)

    # ── Save ──
    out_path = Path(output_dir) / f"RedfireForge-Training-{version}.pptx"
    prs.save(str(out_path))
    print(f"✅ Created: {out_path}")
    return str(out_path)


if __name__ == "__main__":
    version = sys.argv[1] if len(sys.argv) > 1 else "0.2.0"
    script_dir = Path(__file__).parent.parent
    output = script_dir / "training-ppt"
    output.mkdir(exist_ok=True)
    generate(version, str(output))
